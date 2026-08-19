"""parse.py — source slides (.pptx / .pdf) -> IR facts layer.

Two branches, one normalized output:

  * PPTX (high fidelity) via python-pptx. Source groupings are flattened to flat
    geometry (grouping is left to annotators), and known junk placeholder text
    from converted decks is dropped.
  * PDF (medium fidelity) via PyMuPDF. Text blocks, vector drawings, and embedded
    images are extracted. No semantic role inference — a text box is a text box.

All geometry is normalized to [0, 1] of the slide/page. Optional rendering writes
one PNG per slide for visual pairing (best-effort; failures are non-fatal).
"""

from __future__ import annotations

import hashlib
import os
import shutil
import colorsys
import subprocess
import tempfile
from typing import Callable, Dict, Iterable, List, Optional, Tuple

from lxml import etree

from slide_ir import (
    Deck, Slide, Element, Geometry, Style, TextContent, Paragraph, Run, ImageRef,
)

# Junk placeholder strings injected by some "converted" PPTX exporters.
JUNK_TEXT = {"redocEntxTe"}

EMU_PER_PT = 12700


# ===========================================================================
# Public entry points
# ===========================================================================
def load(path: str, deck_id: Optional[str] = None) -> Deck:
    """Parse a .pptx or .pdf into a Deck (facts layer only)."""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".pptx", ".potx"):
        return parse_pptx(path, deck_id=deck_id)
    if ext == ".pdf":
        return parse_pdf(path, deck_id=deck_id)
    raise ValueError(f"unsupported source type: {ext}")


# ===========================================================================
# PPTX branch
# ===========================================================================
def parse_pptx(path: str, deck_id: Optional[str] = None) -> Deck:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(path)
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    theme_cache: Dict[str, Dict[str, str]] = {}
    did = deck_id or _stem(path)
    deck = Deck(deck_id=did, source={
        "origin": "pptx",
        "fidelity": "high",
        "original_filename": os.path.basename(path),
        "parser": "python-pptx",
    })

    for idx, slide in enumerate(prs.slides):
        theme_map = _theme_map_for_slide(slide, theme_cache)
        ctx = _text_color_context(slide, theme_map)
        sid = f"s{idx + 1}"
        s = Slide(id=sid, index=idx,
                  size={"w_emu": sw, "h_emu": sh, "aspect": _aspect(sw, sh)},
                  background={"fill": "#FFFFFF"})
        z = 0
        # Flatten any source groups into absolute-positioned leaf shapes.
        for shape, (l, t, w, h), gfill in _iter_leaves(slide.shapes, identity_xform, theme_map):
            if not _visible(shape, l, t, w, h, sw, sh):
                continue
            el = _pptx_element(shape, l, t, w, h, sw, sh, z, theme_map, gfill, ctx)
            if el is None:
                continue
            if _is_junk(el):
                continue
            s.elements.append(el)
            z += 1
        deck.slides.append(s)
    return deck


def identity_xform(l: float, t: float, w: float, h: float):
    return l, t, w, h


def _iter_leaves(shapes, xform: Callable, theme_map=None, group_fill=None):
    """Yield (leaf_shape, absolute_emu_box, inherited_group_fill_hex)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            child_xform = _group_child_xform(shape, xform)
            gfill = _group_fill_hex(shape, theme_map or {}, group_fill)
            yield from _iter_leaves(shape.shapes, child_xform, theme_map, gfill)
        else:
            try:
                box = xform(int(shape.left), int(shape.top),
                            int(shape.width), int(shape.height))
            except (TypeError, ValueError):
                box = (0, 0, 0, 0)  # shapes occasionally lack a transform
            yield shape, box, group_fill


def _group_fill_hex(group_shape, theme_map, inherited) -> Optional[str]:
    """The fill a group contributes to its <a:grpFill/> children."""
    from pptx.oxml.ns import qn
    grpSpPr = group_shape._element.find(qn("p:grpSpPr"))
    if grpSpPr is None:
        return inherited
    if grpSpPr.find(qn("a:grpFill")) is not None:   # group itself inherits
        return inherited
    if grpSpPr.find(qn("a:noFill")) is not None:
        return None
    sf = grpSpPr.find(qn("a:solidFill"))
    if sf is not None:
        return _color_from_container(sf, theme_map)
    grad = grpSpPr.find(qn("a:gradFill"))
    if grad is not None:
        gslst = grad.find(qn("a:gsLst"))
        gs = gslst.find(qn("a:gs")) if gslst is not None else None
        return _color_from_container(gs, theme_map) if gs is not None else inherited
    return inherited


def _group_child_xform(group_shape, parent_xform: Callable) -> Callable:
    """Compose a coordinate transform for a group's children.

    A group declares its placement (off/ext) and an internal child coordinate
    frame (chOff/chExt). A child's raw coords are in the child frame; map them to
    the group's placed box, then through the parent transform.
    """
    from pptx.oxml.ns import qn

    gl, gt, gw, gh = parent_xform(int(group_shape.left), int(group_shape.top),
                                  int(group_shape.width), int(group_shape.height))
    choff_x = choff_y = 0
    chext_x, chext_y = gw or 1, gh or 1
    try:
        grpSpPr = group_shape._element.find(qn("p:grpSpPr"))
        xfrm = grpSpPr.find(qn("a:xfrm")) if grpSpPr is not None else None
        if xfrm is not None:
            ch_off = xfrm.find(qn("a:chOff"))
            ch_ext = xfrm.find(qn("a:chExt"))
            if ch_off is not None:
                choff_x, choff_y = int(ch_off.get("x")), int(ch_off.get("y"))
            if ch_ext is not None:
                chext_x = int(ch_ext.get("cx")) or 1
                chext_y = int(ch_ext.get("cy")) or 1
    except Exception:
        pass

    sx = gw / chext_x if chext_x else 1.0
    sy = gh / chext_y if chext_y else 1.0

    def xform(l: float, t: float, w: float, h: float):
        return (gl + (l - choff_x) * sx,
                gt + (t - choff_y) * sy,
                w * sx,
                h * sy)

    return xform


def _is_hidden(shape) -> bool:
    """True if any of the shape's non-visual props marks it hidden."""
    from pptx.oxml.ns import qn
    try:
        for cnv in shape._element.iter(qn("p:cNvPr")):
            if cnv.get("hidden") == "1":
                return True
    except Exception:
        pass
    return False


def _visible(shape, l, t, w, h, sw, sh, min_area_frac: float = 1e-6) -> bool:
    """Drop shapes that won't appear on the finished slide: hidden, degenerate,
    fully off-slide, or (optionally) sub-speck area. Lines with one zero
    dimension are kept.
    """
    if _is_hidden(shape):
        return False
    if w < 0 or h < 0:
        return False
    if w == 0 and h == 0:                      # a true point, not a line
        return False
    if l >= sw or t >= sh or (l + w) <= 0 or (t + h) <= 0:   # fully off-slide
        return False
    if w > 0 and h > 0 and min_area_frac > 0 \
            and (w * h) < (sw * sh) * min_area_frac:
        return False
    return True


def _pptx_element(shape, l, t, w, h, sw, sh, z, theme_map=None,
                  group_fill=None, ctx=None) -> Optional[Element]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    etype = _pptx_type(shape)
    if etype is None:
        return None

    geom = Geometry(
        x=_clamp(l / sw), y=_clamp(t / sh),
        w=max(0.0, w / sw), h=max(0.0, h / sh),
        rot=float(getattr(shape, "rotation", 0.0) or 0.0),
    )
    style = _pptx_style(shape, theme_map or {}, group_fill)
    if etype in ("CONNECTOR", "LINE"):
        style.fill = None                          # connectors/lines have no fill
    text = _pptx_text(shape, theme_map or {}, ctx)
    image = None
    if etype == "PICTURE":
        image = _pptx_image(shape)

    eid = f"e{z + 1}"
    return Element(id=eid, type=etype, z=z, geometry=geom,
                   style=style, text=text, image=image)


def _pptx_type(shape) -> Optional[str]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "PICTURE"
    if st == MSO_SHAPE_TYPE.TABLE or getattr(shape, "has_table", False):
        return "TABLE"
    if getattr(shape, "has_chart", False):
        return "CHART"
    if st == MSO_SHAPE_TYPE.LINE:
        return "CONNECTOR"
    if st == MSO_SHAPE_TYPE.FREEFORM:
        return "FREEFORM"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "TEXT_BOX"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "TEXT_BOX" if getattr(shape, "has_text_frame", False) else "AUTO_SHAPE"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return _auto_shape_type(shape)
    # text box drawn as a shape with text
    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
        return "TEXT_BOX"
    return "AUTO_SHAPE"


def _auto_shape_type(shape) -> str:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    try:
        ast = shape.auto_shape_type
    except Exception:
        return "AUTO_SHAPE"
    mapping = {
        MSO_AUTO_SHAPE_TYPE.RECTANGLE: "RECTANGLE",
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE: "ROUNDED_RECTANGLE",
        MSO_AUTO_SHAPE_TYPE.OVAL: "ELLIPSE",
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE: "TRIANGLE",
    }
    return mapping.get(ast, "AUTO_SHAPE")


def _pptx_style(shape, theme_map, group_fill=None) -> Style:
    style = Style()
    style.fill = _shape_fill_hex(shape, theme_map, group_fill)
    style.line_color = _shape_line_hex(shape, theme_map)
    try:
        if shape.line.width is not None:
            style.line_width = round(int(shape.line.width) / EMU_PER_PT, 3)
    except Exception:
        pass
    try:
        if shape.line.dash_style is not None:
            style.line_dash = str(shape.line.dash_style).split()[0].lower()
    except Exception:
        pass
    return style


# ---- theme / scheme color resolution -------------------------------------- #
def _theme_map_for_slide(slide, cache) -> Dict[str, str]:
    """Theme color map for the master this slide actually uses (cached)."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    try:
        master = slide.slide_layout.slide_master
        theme_part = master.part.part_related_by(RT.THEME)
        key = str(theme_part.partname)
        if key not in cache:
            cache[key] = _theme_map_from(theme_part, master)
        return cache[key]
    except Exception:
        return {}


def _theme_map_from(theme_part, master) -> Dict[str, str]:
    from pptx.oxml.ns import qn
    m: Dict[str, str] = {}
    try:
        theme = etree.fromstring(theme_part.blob)
        scheme = theme.find(qn("a:themeElements")).find(qn("a:clrScheme"))
        for child in scheme:                       # dk1, lt1, dk2, lt2, accent1..6, ...
            tag = etree.QName(child).localname
            for c in child:                        # single srgbClr / sysClr child
                hexv = _resolve_clr(c, {})
                if hexv:
                    m[tag] = hexv
                break
        clr_map = master.element.find(qn("p:clrMap"))
        if clr_map is not None:                    # bg1->lt1, tx1->dk1, ...
            for k, v in clr_map.attrib.items():
                if v in m:
                    m[k] = m[v]
    except Exception:
        pass
    return m


def _shape_fill_hex(shape, theme_map, group_fill=None) -> Optional[str]:
    from pptx.oxml.ns import qn
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is not None:
        if spPr.find(qn("a:noFill")) is not None:
            return None
        if spPr.find(qn("a:grpFill")) is not None:   # inherit the enclosing group's fill
            return group_fill
        sf = spPr.find(qn("a:solidFill"))
        if sf is not None:
            return _color_from_container(sf, theme_map)
        grad = spPr.find(qn("a:gradFill"))
        if grad is not None:                       # approximate with first stop
            gslst = grad.find(qn("a:gsLst"))
            gs = gslst.find(qn("a:gs")) if gslst is not None else None
            return _color_from_container(gs, theme_map) if gs is not None else None
        patt = spPr.find(qn("a:pattFill"))
        if patt is not None:                       # approximate with foreground
            fg = patt.find(qn("a:fgClr"))
            return _color_from_container(fg, theme_map) if fg is not None else None
        if spPr.find(qn("a:blipFill")) is not None:
            return None                            # picture fill: no single color
    style = sp.find(qn("p:style"))                 # inherited theme fill
    if style is not None:
        ref = style.find(qn("a:fillRef"))
        if ref is not None:
            if ref.get("idx") in (None, "0"):      # idx 0 == no fill
                return None
            return _color_from_container(ref, theme_map)
    return None


def _shape_line_hex(shape, theme_map) -> Optional[str]:
    from pptx.oxml.ns import qn
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    ln = spPr.find(qn("a:ln")) if spPr is not None else None
    if ln is not None:
        if ln.find(qn("a:noFill")) is not None:
            return None
        sf = ln.find(qn("a:solidFill"))
        if sf is not None:
            return _color_from_container(sf, theme_map)
    style = sp.find(qn("p:style"))
    if style is not None:
        ref = style.find(qn("a:lnRef"))
        if ref is not None:
            if ref.get("idx") in (None, "0"):
                return None
            return _color_from_container(ref, theme_map)
    return None


def _run_color_hex(run, theme_map) -> Optional[str]:
    from pptx.oxml.ns import qn
    try:
        rPr = run._r.find(qn("a:rPr"))
        sf = rPr.find(qn("a:solidFill")) if rPr is not None else None
        if sf is not None:
            return _color_from_container(sf, theme_map)
    except Exception:
        pass
    return _rgb(run.font.color)                    # fallback to explicit rgb


def _color_from_container(container, theme_map) -> Optional[str]:
    from pptx.oxml.ns import qn
    for tag in ("a:srgbClr", "a:schemeClr", "a:sysClr", "a:scrgbClr"):
        c = container.find(qn(tag))
        if c is not None:
            return _resolve_clr(c, theme_map)
    return None


def _resolve_clr(clr, theme_map) -> Optional[str]:
    tag = etree.QName(clr).localname
    if tag == "srgbClr":
        base = "#" + (clr.get("val") or "").upper()
    elif tag == "sysClr":
        last = clr.get("lastClr")
        base = "#" + last.upper() if last else None
    elif tag == "schemeClr":
        base = theme_map.get(clr.get("val"))
    elif tag == "scrgbClr":
        try:
            r = round(int(clr.get("r", "0")) / 100000 * 255)
            g = round(int(clr.get("g", "0")) / 100000 * 255)
            b = round(int(clr.get("b", "0")) / 100000 * 255)
            base = "#%02X%02X%02X" % (r, g, b)
        except Exception:
            base = None
    else:
        base = None
    if not base or len(base) != 7:
        return None
    return _apply_color_mods(base, clr)


def _apply_color_mods(hexv, clr) -> str:
    """Apply lumMod/lumOff/shade/tint children (approximate)."""
    from pptx.oxml.ns import qn
    r, g, b = (int(hexv[i:i + 2], 16) for i in (1, 3, 5))

    def val(tag):
        e = clr.find(qn("a:" + tag))
        return int(e.get("val")) / 100000 if e is not None and e.get("val") else None

    shade, tint = val("shade"), val("tint")
    lm, lo = val("lumMod"), val("lumOff")
    if shade is not None:
        r, g, b = r * shade, g * shade, b * shade
    if tint is not None:
        r = r * tint + 255 * (1 - tint)
        g = g * tint + 255 * (1 - tint)
        b = b * tint + 255 * (1 - tint)
    if lm is not None or lo is not None:
        hh, ll, ss = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
        if lm is not None:
            ll *= lm
        if lo is not None:
            ll += lo
        ll = max(0.0, min(1.0, ll))
        r, g, b = (x * 255 for x in colorsys.hls_to_rgb(hh, ll, ss))
    return "#%02X%02X%02X" % tuple(max(0, min(255, int(round(v)))) for v in (r, g, b))


def _ph_key(sp):
    """(placeholder type, idx) for a shape, or None if not a placeholder."""
    from pptx.oxml.ns import qn
    try:
        ph = sp._element.find(qn("p:nvSpPr")).find(qn("p:nvPr")).find(qn("p:ph"))
    except Exception:
        return None
    if ph is None:
        return None
    return (ph.get("type") or "body", ph.get("idx"))


def _lststyle_color(sp, theme_map) -> Optional[str]:
    """Level-1 default text color from a shape's own <a:lstStyle>, if any."""
    from pptx.oxml.ns import qn
    tb = sp._element.find(qn("p:txBody"))
    lst = tb.find(qn("a:lstStyle")) if tb is not None else None
    lvl = lst.find(qn("a:lvl1pPr")) if lst is not None else None
    defR = lvl.find(qn("a:defRPr")) if lvl is not None else None
    sf = defR.find(qn("a:solidFill")) if defR is not None else None
    return _color_from_container(sf, theme_map) if sf is not None else None


def _txstyles_colors(slide, theme_map) -> Dict[str, str]:
    """Master title/body/other default colors (the lowest-priority fallback)."""
    from pptx.oxml.ns import qn
    out: Dict[str, str] = {}
    try:
        tx = slide.slide_layout.slide_master.element.find(qn("p:txStyles"))
        if tx is None:
            return out
        for key, tag in (("title", "p:titleStyle"),
                         ("body", "p:bodyStyle"),
                         ("other", "p:otherStyle")):
            st = tx.find(qn(tag))
            lvl = st.find(qn("a:lvl1pPr")) if st is not None else None
            defR = lvl.find(qn("a:defRPr")) if lvl is not None else None
            sf = defR.find(qn("a:solidFill")) if defR is not None else None
            if sf is not None:
                c = _color_from_container(sf, theme_map)
                if c:
                    out[key] = c
    except Exception:
        pass
    return out


def _text_color_context(slide, theme_map) -> Dict[str, Any]:
    """Precompute the placeholder text-color inheritance sources for a slide.

    Effective color for an uncolored run resolves in order (most specific first):
    the shape's own list style, the matching LAYOUT placeholder, the matching
    MASTER placeholder, then the master text styles. Layout overrides master.
    """
    ctx = {"ph": {}, "tx": _txstyles_colors(slide, theme_map)}
    try:
        layout = slide.slide_layout
        master = layout.slide_master
        for src in (master, layout):                 # layout overrides master
            for sp in src.shapes:
                if not getattr(sp, "has_text_frame", False):
                    continue
                key = _ph_key(sp)
                if key is None:
                    continue
                col = _lststyle_color(sp, theme_map)
                if col:
                    ctx["ph"][key] = col
                    ctx["ph"][(key[0], None)] = col   # type-only fallback
    except Exception:
        pass
    return ctx


def _inherited_text_color(shape, ctx, theme_map) -> Optional[str]:
    col = _lststyle_color(shape, theme_map)          # shape's own list style wins
    if col:
        return col
    key = _ph_key(shape)
    if key is None:
        return None
    col = ctx["ph"].get(key) or ctx["ph"].get((key[0], None))
    if col:
        return col
    cat = ("title" if key[0] in ("title", "ctrTitle")
           else "body" if key[0] in ("body", "subTitle") else "other")
    return ctx["tx"].get(cat) or ctx["tx"].get("other")


def _pptx_text(shape, theme_map, ctx=None) -> Optional[TextContent]:
    if not getattr(shape, "has_text_frame", False):
        return None
    inherited = _inherited_text_color(shape, ctx or {"ph": {}, "tx": {}}, theme_map)

    paras: List[Paragraph] = []
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            f = r.font
            runs.append(Run(
                t=r.text,
                font=f.name,
                size_pt=(f.size.pt if f.size is not None else None),
                bold=f.bold,
                italic=f.italic,
                color=_run_color_hex(r, theme_map) or inherited,   # WYSIWYG effective color
            ))
        if not runs and p.text:
            runs = [Run(t=p.text)]
        align = (str(p.alignment).split()[0].lower()
                 if p.alignment is not None else "left")
        paras.append(Paragraph(align=align, runs=runs))
    tc = TextContent(paras)
    return tc if tc.plain().strip() else None


def _pptx_image(shape) -> Optional[ImageRef]:
    try:
        blob = shape.image.blob
        return ImageRef(hash="sha256:" + hashlib.sha256(blob).hexdigest(),
                        ext=shape.image.ext,
                        alt=(shape.name or None))
    except Exception:
        return None


def _rgb(color_obj) -> Optional[str]:
    """Return '#RRGGBB' for an explicit RGB color, else None (theme colors)."""
    try:
        rgb = color_obj.rgb
        return "#" + str(rgb)
    except Exception:
        return None


# ===========================================================================
# PDF branch
# ===========================================================================
def parse_pdf(path: str, deck_id: Optional[str] = None) -> Deck:
    import fitz  # PyMuPDF; install with: pip install pymupdf

    doc = fitz.open(path)
    did = deck_id or _stem(path)
    deck = Deck(deck_id=did, source={
        "origin": "pdf",
        "fidelity": "medium",
        "original_filename": os.path.basename(path),
        "parser": "PyMuPDF",
    })

    for idx, page in enumerate(doc):
        pw, ph = page.rect.width, page.rect.height
        sid = f"s{idx + 1}"
        s = Slide(id=sid, index=idx,
                  size={"w_pt": pw, "h_pt": ph, "aspect": _aspect(pw, ph)},
                  background={"fill": "#FFFFFF"})
        z = 0

        # 1) vector drawings -> shapes
        for d in page.get_drawings():
            el = _pdf_drawing(d, pw, ph, z)
            if el is not None:
                s.elements.append(el)
                z += 1

        # 2) text blocks -> text boxes
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type", 0) != 0:      # 0 == text block
                continue
            el = _pdf_text_block(block, pw, ph, z)
            if el is None or _is_junk(el):
                continue
            s.elements.append(el)
            z += 1

        # 3) embedded images -> pictures
        for img in page.get_images(full=True):
            el = _pdf_image(doc, page, img, pw, ph, z)
            if el is not None:
                s.elements.append(el)
                z += 1

        deck.slides.append(s)
    return deck


def _pdf_drawing(d: dict, pw: float, ph: float, z: int) -> Optional[Element]:
    rect = d.get("rect")
    if rect is None:
        return None
    x0, y0, x1, y1 = rect
    geom = Geometry(x=_clamp(x0 / pw), y=_clamp(y0 / ph),
                    w=max(0.0, (x1 - x0) / pw), h=max(0.0, (y1 - y0) / ph))
    items = d.get("items", [])
    kinds = {it[0] for it in items}
    if kinds == {"l"}:
        etype = "LINE"
    elif kinds <= {"re"}:
        etype = "RECTANGLE"
    else:
        etype = "FREEFORM"
    style = Style(
        fill=_pdf_color(d.get("fill")),
        line_color=_pdf_color(d.get("color")),
        line_width=(round(d["width"], 3) if d.get("width") else None),
    )
    return Element(id=f"e{z + 1}", type=etype, z=z, geometry=geom, style=style)


def _pdf_text_block(block: dict, pw: float, ph: float, z: int) -> Optional[Element]:
    x0, y0, x1, y1 = block["bbox"]
    geom = Geometry(x=_clamp(x0 / pw), y=_clamp(y0 / ph),
                    w=max(0.0, (x1 - x0) / pw), h=max(0.0, (y1 - y0) / ph))
    paras: List[Paragraph] = []
    for line in block.get("lines", []):
        runs = []
        for sp in line.get("spans", []):
            flags = sp.get("flags", 0)
            runs.append(Run(
                t=sp.get("text", ""),
                font=sp.get("font"),
                size_pt=round(sp.get("size", 0.0), 2) or None,
                bold=bool(flags & 16),
                italic=bool(flags & 2),
                color=_pdf_color(sp.get("color")),
            ))
        if runs:
            paras.append(Paragraph(align="left", runs=runs))
    tc = TextContent(paras)
    if not tc.plain().strip():
        return None
    return Element(id=f"e{z + 1}", type="TEXT_BOX", z=z, geometry=geom, text=tc)


def _pdf_image(doc, page, img, pw: float, ph: float, z: int) -> Optional[Element]:
    xref = img[0]
    try:
        rects = page.get_image_rects(xref)
    except Exception:
        rects = []
    if not rects:
        return None
    r = rects[0]
    geom = Geometry(x=_clamp(r.x0 / pw), y=_clamp(r.y0 / ph),
                    w=max(0.0, (r.x1 - r.x0) / pw), h=max(0.0, (r.y1 - r.y0) / ph))
    try:
        info = doc.extract_image(xref)
        h = "sha256:" + hashlib.sha256(info["image"]).hexdigest()
        ext = info.get("ext")
    except Exception:
        h, ext = f"xref:{xref}", None
    return Element(id=f"e{z + 1}", type="PICTURE", z=z, geometry=geom,
                   image=ImageRef(hash=h, ext=ext))


def _pdf_color(c) -> Optional[str]:
    """PyMuPDF colors are ints (sRGB) or float tuples in [0,1]; -> '#RRGGBB'."""
    if c is None:
        return None
    if isinstance(c, int):
        return "#%06X" % (c & 0xFFFFFF)
    if isinstance(c, (tuple, list)) and len(c) >= 3:
        r, g, b = (int(round(v * 255)) for v in c[:3])
        return "#%02X%02X%02X" % (r, g, b)
    return None


# ===========================================================================
# Rendering (best-effort visual pairing)
# ===========================================================================
def render_slides(src_path: str, deck: Deck, out_dir: str, dpi: int = 150) -> None:
    """Write one PNG per slide into out_dir/images and set slide.image.

    PDF rendering uses PyMuPDF directly. PPTX rendering converts to PDF via
    LibreOffice (`soffice`) first. Any failure is logged and left non-fatal.
    """
    img_dir = os.path.join(out_dir, "images")
    os.makedirs(img_dir, exist_ok=True)
    ext = os.path.splitext(src_path)[1].lower()
    pdf_path = src_path

    try:
        if ext in (".pptx", ".potx"):
            pdf_path = _soffice_to_pdf(src_path, out_dir)
            if pdf_path is None:
                print("[render] soffice not available; skipping PPTX render.")
                return
        _render_pdf_pages(pdf_path, deck, img_dir, dpi)
    except Exception as exc:  # rendering is optional
        print(f"[render] skipped ({exc})")


def _render_pdf_pages(pdf_path: str, deck: Deck, img_dir: str, dpi: int) -> None:
    import fitz
    doc = fitz.open(pdf_path)
    for idx, page in enumerate(doc):
        if idx >= len(deck.slides):
            break
        png = os.path.join(img_dir, f"{deck.deck_id}_s{idx + 1}.png")
        page.get_pixmap(dpi=dpi).save(png)
        deck.slides[idx].image = os.path.relpath(png, os.path.dirname(img_dir))


def _soffice_to_pdf(src_path: str, out_dir: str) -> Optional[str]:
    soffice = (shutil.which("soffice") or shutil.which("libreoffice")
               or "/Applications/LibreOffice.app/Contents/MacOS/soffice")
    if not soffice or not os.path.exists(soffice):
        return None
    tmp = tempfile.mkdtemp()
    out_pdf = os.path.join(tmp, _stem(src_path) + ".pdf")
    # Include hidden ("skip") slides so PDF pages align 1:1 with parsed slides.
    # (LibreOffice drops them by default, which shifts every later image.)
    filter_arg = ('pdf:impress_pdf_Export:'
                  '{"ExportHiddenSlides":{"type":"boolean","value":"true"}}')
    for arg in (filter_arg, "pdf"):          # fall back if filter syntax unsupported
        try:
            subprocess.run([soffice, "--headless", "--convert-to", arg,
                            "--outdir", tmp, src_path],
                           check=True, capture_output=True, timeout=120)
        except Exception:
            pass
        if os.path.exists(out_pdf):
            return out_pdf
    return None


# ===========================================================================
# small helpers
# ===========================================================================
def _is_junk(el: Element) -> bool:
    return bool(el.text) and el.text.plain().strip() in JUNK_TEXT


def _clamp(v: float) -> float:
    return max(0.0, min(1.0, v))


def _aspect(w: float, h: float) -> str:
    if not w or not h:
        return "?"
    r = w / h
    return "16:9" if abs(r - 16 / 9) < 0.05 else ("4:3" if abs(r - 4 / 3) < 0.05
                                                  else f"{r:.2f}:1")


def _stem(path: str) -> str:
    return os.path.splitext(os.path.basename(path))[0]